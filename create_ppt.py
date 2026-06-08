from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x00, 0x78, 0xFF)
DARK = RGBColor(0x04, 0x0B, 0x14)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x4E, 0x4E, 0x4E)
LIGHT_GRAY = RGBColor(0x6C, 0x75, 0x7D)

def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color, left=0, top=0, width=None, height=None):
    if width is None: width = prs.slide_width
    if height is None: height = prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_text(slide, items, left, top, width, height, size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸ {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return txBox

# --- SLIDE 1: TITLE ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape_bg(slide, BLUE, Inches(0), Inches(2.8), Inches(0.15), Inches(2))
add_text(slide, "NASIR HAROON", Inches(1.5), Inches(2.8), Inches(10), Inches(1), size=48, bold=True, color=WHITE)
add_text(slide, "Full Stack Developer", Inches(1.5), Inches(3.8), Inches(10), Inches(0.6), size=28, color=RGBColor(0x00, 0xA6, 0xFF))
add_text(slide, "PHP | Laravel | Symfony | FastAPI | React", Inches(1.5), Inches(4.5), Inches(10), Inches(0.5), size=18, color=LIGHT_GRAY)
add_text(slide, "5+ Years Experience  |  30+ Projects  |  15+ Clients", Inches(1.5), Inches(5.3), Inches(10), Inches(0.5), size=16, color=LIGHT_GRAY)

# --- SLIDE 2: ABOUT ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, BLUE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
add_text(slide, "ABOUT ME", Inches(0.8), Inches(0.2), Inches(8), Inches(0.8), size=32, bold=True, color=WHITE)
add_text(slide, "Full Stack Developer | Islamabad, Pakistan", Inches(0.8), Inches(0.7), Inches(8), Inches(0.5), size=16, color=RGBColor(0xCC, 0xE5, 0xFF))

items = [
    "5+ years building scalable web apps with PHP, Laravel, Symfony, React, Vue, and FastAPI",
    "Expert in designing RESTful APIs, optimizing databases, and secure architectures",
    "Built crawlers for 50+ European property sources, reducing manual work by 80%",
    "Mentored teams and delivered projects serving thousands of users across Belgium, France, and Spain",
    "AI-augmented development workflows including MCP agent integration and spec-driven coding",
]
add_bullet_text(slide, items, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5), size=18, color=DARK)

# Contact info
add_text(slide, "✉ sardarnasirharoon@gmail.com", Inches(0.8), Inches(6.2), Inches(5), Inches(0.4), size=14, color=BLUE)
add_text(slide, "📞 +92 346 4382740", Inches(6.5), Inches(6.2), Inches(5), Inches(0.4), size=14, color=BLUE)

# --- SLIDE 3: SKILLS ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape_bg(slide, BLUE, Inches(0), Inches(0), Inches(0.15), Inches(2))
add_text(slide, "TECHNICAL SKILLS", Inches(1.5), Inches(0.5), Inches(10), Inches(1), size=36, bold=True, color=WHITE)

skills = [
    ("PHP / Laravel / Symfony / CodeIgniter", "95%"),
    ("JavaScript / React / Vue / TypeScript", "85%"),
    ("Python / FastAPI / Scrapy / BeautifulSoup", "80%"),
    ("PostgreSQL / MySQL / MongoDB", "90%"),
    ("RESTful API Design & Integration", "90%"),
    ("Docker / Git / CI-CD / Jira", "75%"),
    ("Web Scraping & Automation", "85%"),
    ("Team Leadership & Code Review", "80%"),
]

for i, (skill, pct) in enumerate(skills):
    y = Inches(1.8) + Inches(i * 0.65)
    add_text(slide, skill, Inches(1.5), y, Inches(7), Inches(0.4), size=16, color=WHITE)
    add_text(slide, pct, Inches(11), y, Inches(1.5), Inches(0.4), size=16, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
    # progress bar bg
    bar_y = y + Inches(0.35)
    add_shape_bg(slide, RGBColor(0x2A, 0x2A, 0x3A), Inches(1.5), bar_y, Inches(10.5), Inches(0.12))
    # progress bar fill - parse percentage
    val = int(pct.replace("%", "")) / 100
    add_shape_bg(slide, BLUE, Inches(1.5), bar_y, Inches(10.5 * val), Inches(0.12))

# --- SLIDE 4: EXPERIENCE ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, BLUE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
add_text(slide, "PROFESSIONAL EXPERIENCE", Inches(0.8), Inches(0.2), Inches(10), Inches(0.8), size=32, bold=True, color=WHITE)

exps = [
    ("Software Engineer - Full Stack", "Zeropoint IT, Islamabad", "Feb 2023 - Present"),
    ("Software Engineer - Laravel", "Syslab Technologies, Malaysia", "May 2023 - Feb 2024"),
    ("Software Engineer - PHP", "MicroMerger (Pvt.) Ltd.", "Nov 2021 - Jun 2023"),
    ("Associate Software Engineer", "Pace Technologies (Pvt.) Ltd.", "Jan 2021 - Nov 2021"),
]

for i, (title, company, period) in enumerate(exps):
    y = Inches(1.8) + Inches(i * 1.3)
    # blue dot + line
    add_shape_bg(slide, BLUE, Inches(1.2), y + Inches(0.15), Inches(0.12), Inches(0.12))
    if i < len(exps) - 1:
        add_shape_bg(slide, RGBColor(0xDD, 0xDD, 0xDD), Inches(1.25), y + Inches(0.3), Inches(0.04), Inches(1))
    add_text(slide, title, Inches(1.8), y, Inches(7), Inches(0.4), size=18, bold=True, color=DARK)
    add_text(slide, company, Inches(1.8), y + Inches(0.4), Inches(7), Inches(0.3), size=14, color=GRAY)
    add_text(slide, period, Inches(1.8), y + Inches(0.7), Inches(4), Inches(0.3), size=12, color=BLUE)

# --- SLIDE 5: PORTFOLIO OVERVIEW ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape_bg(slide, BLUE, Inches(0), Inches(0), Inches(0.15), Inches(2))
add_text(slide, "KEY PROJECTS", Inches(1.5), Inches(0.5), Inches(10), Inches(1), size=36, bold=True, color=WHITE)

projects_data = [
    ("SDG3 Pakistan", "Health data platform, UN SDG Goal 3 — CodeIgniter + PostgreSQL"),
    ("EPI - MIS Federal", "Centralized immunization MIS for Pakistan — PHP, CodeIgniter"),
    ("EPI - MIS Punjab/Sindh/Balochistan", "Provincial immunization systems for 3 provinces"),
    ("EPI - MIS KPK/GB/AJK/Islamabad", "4 additional provincial immunization platforms"),
    ("XpertJobs", "Job portal connecting top employers & professionals — Laravel"),
]

for i, (name, desc) in enumerate(projects_data):
    y = Inches(1.8) + Inches(i * 1.05)
    add_text(slide, f"0{i+1}", Inches(1.5), y, Inches(0.8), Inches(0.6), size=28, bold=True, color=BLUE)
    add_text(slide, name, Inches(2.5), y, Inches(4), Inches(0.4), size=18, bold=True, color=WHITE)
    add_text(slide, desc, Inches(2.5), y + Inches(0.4), Inches(9), Inches(0.4), size=14, color=LIGHT_GRAY)

# --- SLIDE 6: ACHIEVEMENTS ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, BLUE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
add_text(slide, "ACHIEVEMENTS & HIGHLIGHTS", Inches(0.8), Inches(0.2), Inches(10), Inches(0.8), size=32, bold=True, color=WHITE)

stats_data = [
    ("15+", "Clients Served"),
    ("30+", "Projects Delivered"),
    ("5+", "Years Experience"),
    ("15+", "Technologies"),
    ("99.9%", "Uptime"),
    ("50+", "Data Sources"),
]

for i, (num, label) in enumerate(stats_data):
    col = i % 3
    row = i // 3
    x = Inches(1.5) + Inches(col * 4)
    y = Inches(1.8) + Inches(row * 2.2)
    shape = add_shape_bg(slide, BLUE, x, y, Inches(3), Inches(1.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    add_text(slide, num, x, y + Inches(0.3), Inches(3), Inches(0.7), size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + Inches(1.1), Inches(3), Inches(0.4), size=16, color=RGBColor(0xCC, 0xE5, 0xFF), align=PP_ALIGN.CENTER)

# --- SLIDE 7: CONTACT ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape_bg(slide, BLUE, Inches(0), Inches(3), Inches(13.333), Inches(0.08))
add_text(slide, "LET'S WORK TOGETHER", Inches(1.5), Inches(1.5), Inches(10), Inches(1), size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "I'm available for freelance projects, full-time roles, and collaborations", Inches(1.5), Inches(2.5), Inches(10), Inches(0.6), size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, "✉ sardarnasirharoon@gmail.com", Inches(1.5), Inches(4), Inches(5), Inches(0.5), size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "📞 +92 346 4382740", Inches(6.5), Inches(4), Inches(5), Inches(0.5), size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "🔗 linkedin.com/in/nasirharoon786", Inches(1.5), Inches(4.6), Inches(5), Inches(0.5), size=18, color=BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "🐙 github.com/nasirharoon786", Inches(6.5), Inches(4.6), Inches(5), Inches(0.5), size=18, color=BLUE, align=PP_ALIGN.CENTER)

add_text(slide, "nasir-haroon.vercel.app", Inches(1.5), Inches(5.8), Inches(10), Inches(0.5), size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Save
path = "/home/nasir/Downloads/DevFolio/Nasir_Haroon_Portfolio.pptx"
os.makedirs(os.path.dirname(path), exist_ok=True)
prs.save(path)
print(f"Presentation saved to: {path}")
